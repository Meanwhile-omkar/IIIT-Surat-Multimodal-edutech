"""Handles parsing of PDF, PPTX, and TXT files."""

import os
import tempfile
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation

from app.core.config import UPLOAD_DIR


def parse_pdf(file_path: str) -> str:
    """Extract text from PDF using PyMuPDF."""
    doc = fitz.open(file_path)
    text_parts = []
    for page in doc:
        text_parts.append(page.get_text())
    doc.close()
    return "\n\n".join(text_parts)


def parse_pptx(file_path: str) -> str:
    """Extract text from PPTX, preserving slide order."""
    prs = Presentation(file_path)
    text_parts = []
    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        slide_texts.append(line)
        if slide_texts:
            text_parts.append(f"[Slide {i+1}]\n" + "\n".join(slide_texts))
    return "\n\n".join(text_parts)


def parse_txt(file_path: str) -> str:
    """Read plain text file."""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def save_upload(file_bytes: bytes, filename: str) -> str:
    """Save uploaded file to disk. Returns file path."""
    os.makedirs(UPLOAD_DIR, exist_ok=True)
    file_path = os.path.join(UPLOAD_DIR, filename)
    with open(file_path, "wb") as f:
        f.write(file_bytes)
    return file_path


def parse_document(file_path: str) -> tuple[str, str]:
    """
    Parse a document based on its extension.
    Returns (extracted_text, file_type).
    """
    ext = Path(file_path).suffix.lower()

    if ext == ".pdf":
        return parse_pdf(file_path), "pdf"
    elif ext == ".pptx":
        return parse_pptx(file_path), "pptx"
    elif ext in (".txt", ".md"):
        return parse_txt(file_path), "txt"
    else:
        raise ValueError(f"Unsupported file type: {ext}. Supported: .pdf, .pptx, .txt, .md")
